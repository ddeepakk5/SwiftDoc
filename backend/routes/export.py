from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from database import get_db
from models import Project, Section
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.xmlchemy import OxmlElement
from docx.shared import Pt as DocxPt 
from docx.shared import Inches # <--- NEW IMPORT FOR INDENTATION
import os

router = APIRouter()

# --- HELPER: Turn off bullets AND Fix Indentation (PPTX) ---
def remove_bullet_formatting(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    if pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone") is None:
        buNone = OxmlElement('a:buNone')
        pPr.insert(0, buNone)
    pPr.set('marL', '0')
    pPr.set('indent', '0')

# --- HELPER: Parse Markdown for Word (.docx) ---
def add_markdown_to_docx(doc, text):
    lines = text.split('\n')
    for line in lines:
        # We need the raw line first to count spaces
        if not line.strip():
            continue
            
        # 1. Calculate Indentation Level
        # Count leading spaces. We assume 2 spaces = 1 indent level.
        leading_spaces = len(line) - len(line.lstrip())
        indent_level = leading_spaces // 2 
        
        stripped_line = line.strip()

        # 2. Detect Bullet Points vs Paragraphs
        if stripped_line.startswith('* ') or stripped_line.startswith('- '):
            style = 'List Bullet'
            clean_line = stripped_line[2:]
            
            p = doc.add_paragraph(style=style)
            p.paragraph_format.space_after = DocxPt(0) 
            
            # --- APPLY NESTED INDENTATION ---
            if indent_level > 0:
                # Standard indent is usually 0.25 inches. 
                # We multiply by (indent_level + 1) to push sub-points right.
                p.paragraph_format.left_indent = Inches(0.25 * (indent_level + 1))
                
        else:
            style = 'Normal'
            clean_line = stripped_line
            p = doc.add_paragraph(style=style)
            p.paragraph_format.space_after = DocxPt(12) 

        # 3. Detect Formatting: Bold (**) AND Italics (_)
        bold_parts = clean_line.split('**')
        for i, bold_part in enumerate(bold_parts):
            is_bold = (i % 2 == 1) 
            
            italic_parts = bold_part.split('_')
            for j, sub_text in enumerate(italic_parts):
                if not sub_text: continue 
                
                is_italic = (j % 2 == 1) 
                
                run = p.add_run(sub_text)
                run.bold = is_bold
                run.italic = is_italic

# --- HELPER: Parse Markdown for PowerPoint (.pptx) ---
def add_markdown_to_pptx(text_frame, text):
    text_frame.clear() 
    p = text_frame.paragraphs[0]
    
    lines = text.split('\n')
    valid_lines = [l for l in lines if l.strip()] # Keep raw line to check indent
    
    for idx, line in enumerate(valid_lines):
        if idx > 0:
            p = text_frame.add_paragraph()

        stripped_line = line.strip()
        
        # Calculate Indentation for PPT
        leading_spaces = len(line) - len(line.lstrip())
        indent_level = leading_spaces // 2

        # Bullet Points
        if stripped_line.startswith('* ') or stripped_line.startswith('- '):
            clean_line = stripped_line[2:]
            
            # PPT supports levels 0-8 directly
            p.level = min(indent_level, 8) 
        
        # Normal Paragraphs
        else:
            clean_line = stripped_line
            p.level = 0 
            remove_bullet_formatting(p)
            p.space_after = Pt(14)

        # Bold & Formatting
        parts = clean_line.split('**')
        for i, part in enumerate(parts):
            run = p.add_run()
            run.text = part
            run.font.size = Pt(18)
            if i % 2 == 1:
                run.font.bold = True

@router.get("/projects/{project_id}/export")
def export_project(project_id: int, db: Session = Depends(get_db)):
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(404, "Project not found")
    
    sections = db.query(Section).filter(Section.project_id == project_id).order_by(Section.order_index).all()
    filename = f"export_{project_id}.{project.doc_type}"
    
    if os.path.exists(filename):
        os.remove(filename)

    if project.doc_type == 'docx':
        doc = Document()
        doc.add_heading(project.title, 0)
        
        p = doc.add_paragraph(project.topic)
        p.style = 'Subtitle'
        
        for sec in sections:
            doc.add_heading(sec.title, level=1)
            if sec.content:
                add_markdown_to_docx(doc, sec.content)
            else:
                doc.add_paragraph("[No content]")
        doc.save(filename)
        
    elif project.doc_type == 'pptx':
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = project.title
        slide.placeholders[1].text = project.topic
        
        for sec in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sec.title
            content_placeholder = slide.placeholders[1]
            if sec.content:
                add_markdown_to_pptx(content_placeholder.text_frame, sec.content)
            else:
                content_placeholder.text = "[No content]"
        prs.save(filename)
    
    return FileResponse(filename, filename=filename)